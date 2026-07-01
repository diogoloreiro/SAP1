#!/usr/bin/env python3
# =============================================================
# gerar_pptx.py - Apresentacao PROFISSIONAL da arquitetura do SAP-1
# Layouts variados (cards, callouts, codigo, tabelas, imagens),
# 3 apresentadores, notas do apresentador, rodape com numeracao.
# Uso (na pasta apresentacao):  python3 gerar_pptx.py
# =============================================================
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

DISCIPLINA = "Arquitetura de Computadores"      # <-- ajuste
INSTITUICAO = "[Instituição] · [Curso]"          # <-- ajuste
DATA = "2025"                                     # <-- ajuste

DARK  = RGBColor(0x0D,0x2B,0x4E); MID = RGBColor(0x0D,0x3B,0x66)
BLUE  = RGBColor(0x15,0x65,0xC0); TEAL= RGBColor(0x00,0x83,0x8F)
GREY  = RGBColor(0x5A,0x63,0x6E); LGREY=RGBColor(0xEC,0xF1,0xF5)
WHITE = RGBColor(0xFF,0xFF,0xFF); INK = RGBColor(0x1B,0x26,0x31)
GREEN = RGBColor(0x2E,0x7D,0x32); ORANGE=RGBColor(0xE6,0x5C,0x00)
CODEBG= RGBColor(0x1E,0x24,0x30); CODEFG=RGBColor(0xE6,0xED,0xF3)
PRES  = {"Pessoa 1": BLUE, "Pessoa 2": TEAL, "Pessoa 3": ORANGE}

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height
META = []   # (slide, presenter or None)

def _tb(s,l,t,w,h):
    tf=s.shapes.add_textbox(l,t,w,h).text_frame; tf.word_wrap=True
    tf.margin_left=Pt(2); tf.margin_right=Pt(2); return tf
def _rect(s,l,t,w,h,c,rounded=False,line=None):
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,l,t,w,h)
    shp.fill.solid(); shp.fill.fore_color.rgb=c
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb=line; shp.line.width=Pt(1)
    shp.shadow.inherit=False; return shp
def _para(tf,text,size,color,bold=False,italic=False,align=None,after=6,first=False,level=0):
    p=tf.paragraphs[0] if first else tf.add_paragraph()
    p.text=text; p.font.size=Pt(size); p.font.color.rgb=color
    p.font.bold=bold; p.font.italic=italic; p.level=level
    if align is not None: p.alignment=align
    p.space_after=Pt(after); return p
def notes(s,t): s.notes_slide.notes_text_frame.text=t

def _header(s,title,cor,kicker=None):
    _rect(s,0,0,SW,Inches(1.2),MID)
    _rect(s,0,Inches(1.2),SW,Inches(0.07),cor)
    if kicker:
        tf=_tb(s,Inches(0.62),Inches(0.16),Inches(11),Inches(0.35))
        _para(tf,kicker.upper(),12,RGBColor(0x8F,0xB8,0xDE),bold=True,after=0,first=True)
        tf2=_tb(s,Inches(0.6),Inches(0.5),Inches(12.1),Inches(0.62))
        _para(tf2,title,26,WHITE,bold=True,after=0,first=True)
    else:
        tf=_tb(s,Inches(0.6),Inches(0.2),Inches(12.1),Inches(0.8)); tf.vertical_anchor=MSO_ANCHOR.MIDDLE
        _para(tf,title,28,WHITE,bold=True,after=0,first=True)

# ---------- slides ----------
def cover(title,sub):
    s=prs.slides.add_slide(BLANK); _rect(s,0,0,SW,SH,DARK)
    _rect(s,0,0,Inches(0.28),SH,ORANGE)
    tf=_tb(s,Inches(0.95),Inches(0.85),Inches(11.4),Inches(0.4))
    _para(tf,DISCIPLINA.upper(),14,RGBColor(0x8F,0xB8,0xDE),bold=True,first=True)
    tf=_tb(s,Inches(0.9),Inches(1.7),Inches(11.6),Inches(1.9))
    _para(tf,title,48,WHITE,bold=True,first=True,after=4)
    _para(tf,sub,20,RGBColor(0xC7,0xDD,0xF0))
    _rect(s,Inches(0.95),Inches(4.35),Inches(4.2),Inches(0.045),ORANGE)
    tf=_tb(s,Inches(0.92),Inches(4.6),Inches(11.5),Inches(1.7))
    _para(tf,"Integrante 1  ·  [nome]",18,WHITE,first=True,after=4)
    _para(tf,"Integrante 2  ·  [nome]",18,WHITE,after=4)
    _para(tf,"Integrante 3  ·  [nome]",18,WHITE,after=4)
    tf=_tb(s,Inches(0.92),Inches(6.6),Inches(11.5),Inches(0.5))
    _para(tf,f"{INSTITUICAO}   ·   {DATA}",14,RGBColor(0x9F,0xB6,0xCC),first=True)
    META.append((s,None)); return s

def section(num,tit,ap,topics):
    s=prs.slides.add_slide(BLANK); cor=PRES[ap]; _rect(s,0,0,SW,SH,cor)
    _rect(s,0,0,SW,SH,cor)
    tf=_tb(s,Inches(0.95),Inches(1.9),Inches(11),Inches(0.5))
    _para(tf,f"PARTE {num} DE 3",16,WHITE,bold=True,first=True)
    tf=_tb(s,Inches(0.9),Inches(2.5),Inches(11.4),Inches(1.3))
    _para(tf,tit,42,WHITE,bold=True,first=True)
    _rect(s,Inches(0.95),Inches(3.95),Inches(3.6),Inches(0.05),WHITE)
    tf=_tb(s,Inches(0.95),Inches(4.2),Inches(11),Inches(2.2))
    for i,t in enumerate(topics):
        _para(tf,"•  "+t,18,WHITE,first=(i==0),after=7)
    tf=_tb(s,Inches(0.95),Inches(6.5),Inches(11),Inches(0.5))
    _para(tf,"Apresenta: "+ap,17,RGBColor(0xEE,0xEE,0xEE),bold=True,first=True)
    META.append((s,None)); return s

def agenda(parts):
    s=prs.slides.add_slide(BLANK); _header(s,"Sumário",BLUE)
    x=Inches(0.7); w=Inches(3.95); gap=Inches(0.28); top=Inches(1.75)
    for i,(cor,titulo,ap,itens) in enumerate(parts):
        lx=x+ (w+gap)*i
        card=_rect(s,lx,top,w,Inches(4.9),LGREY,rounded=True)
        _rect(s,lx,top,w,Inches(0.9),cor,rounded=True)
        tf=_tb(s,lx+Inches(0.25),top+Inches(0.12),w-Inches(0.5),Inches(0.7))
        _para(tf,f"PARTE {i+1}",12,WHITE,bold=True,first=True,after=0)
        _para(tf,titulo,17,WHITE,bold=True,after=0)
        tf=_tb(s,lx+Inches(0.28),top+Inches(1.15),w-Inches(0.5),Inches(3.0))
        for j,it in enumerate(itens):
            _para(tf,"•  "+it,13.5,INK,first=(j==0),after=6)
        tf=_tb(s,lx+Inches(0.28),top+Inches(4.35),w-Inches(0.5),Inches(0.45))
        _para(tf,ap,13,cor,bold=True,first=True)
    META.append((s,"Todos")); return s

def content(title,bullets,ap,nota,kicker=None,takeaway=None):
    s=prs.slides.add_slide(BLANK); _header(s,title,PRES[ap],kicker)
    bh=Inches(4.6) if takeaway else Inches(5.1)
    body=_tb(s,Inches(0.75),Inches(1.55),Inches(11.8),bh)
    for i,it in enumerate(bullets):
        txt,lvl=(it if isinstance(it,tuple) else (it,0))
        _para(body,("▸  " if lvl==0 else "–  ")+txt, 20 if lvl==0 else 17,
              INK if lvl==0 else GREY, first=(i==0), after=10, level=lvl)
    if takeaway:
        box=_rect(s,Inches(0.75),Inches(6.15),Inches(11.85),Inches(0.72),LGREY,rounded=True)
        _rect(s,Inches(0.75),Inches(6.15),Inches(0.12),Inches(0.72),PRES[ap])
        tf=_tb(s,Inches(1.05),Inches(6.2),Inches(11.4),Inches(0.62)); tf.vertical_anchor=MSO_ANCHOR.MIDDLE
        _para(tf,"Chave:  "+takeaway,15,MID,bold=True,first=True,after=0)
    META.append((s,ap)); notes(s,nota); return s

def cards(title,items,ap,nota,kicker=None):
    s=prs.slides.add_slide(BLANK); _header(s,title,PRES[ap],kicker)
    n=len(items); cols=2 if n<=4 else 3
    import math; rows=math.ceil(n/cols)
    x0=Inches(0.7); y0=Inches(1.65)
    cw=(SW-Inches(1.4)-Inches(0.3)*(cols-1))/cols
    ch=(Inches(5.15)-Inches(0.3)*(rows-1))/rows
    for i,(t,body) in enumerate(items):
        r,c=divmod(i,cols)
        lx=x0+ (cw+Inches(0.3))*c; ty=y0+(ch+Inches(0.3))*r
        _rect(s,lx,ty,cw,ch,LGREY,rounded=True)
        _rect(s,lx,ty,cw,Inches(0.12),PRES[ap])
        tf=_tb(s,lx+Inches(0.22),ty+Inches(0.2),cw-Inches(0.44),ch-Inches(0.4))
        _para(tf,t,17,MID,bold=True,first=True,after=5)
        _para(tf,body,14,GREY,after=0)
    META.append((s,ap)); notes(s,nota); return s

def code_slide(title,intro,code,ap,nota,takeaway=None,kicker=None):
    s=prs.slides.add_slide(BLANK); _header(s,title,PRES[ap],kicker)
    tf=_tb(s,Inches(0.75),Inches(1.5),Inches(11.8),Inches(0.7))
    _para(tf,intro,18,INK,first=True,after=0)
    _rect(s,Inches(0.75),Inches(2.35),Inches(11.85),Inches(3.4),CODEBG,rounded=True)
    tf=_tb(s,Inches(1.05),Inches(2.55),Inches(11.3),Inches(3.0))
    for i,ln in enumerate(code.split("\n")):
        p=_para(tf,ln if ln else " ",15,CODEFG,first=(i==0),after=2)
        p.font.name="Consolas"
    if takeaway:
        box=_rect(s,Inches(0.75),Inches(6.0),Inches(11.85),Inches(0.72),LGREY,rounded=True)
        _rect(s,Inches(0.75),Inches(6.0),Inches(0.12),Inches(0.72),PRES[ap])
        tf=_tb(s,Inches(1.05),Inches(6.05),Inches(11.4),Inches(0.62)); tf.vertical_anchor=MSO_ANCHOR.MIDDLE
        _para(tf,"Chave:  "+takeaway,15,MID,bold=True,first=True,after=0)
    META.append((s,ap)); notes(s,nota); return s

def image(title,img,cap,ap,nota,kicker=None):
    s=prs.slides.add_slide(BLANK); _header(s,title,PRES[ap],kicker)
    pic=s.shapes.add_picture(img,0,Inches(1.5),height=Inches(4.95))
    pic.left=int((SW-pic.width)/2)
    if cap:
        tf=_tb(s,Inches(0.6),Inches(6.55),Inches(12.1),Inches(0.4))
        _para(tf,cap,13,GREY,italic=True,align=PP_ALIGN.CENTER,first=True,after=0)
    META.append((s,ap)); notes(s,nota); return s

def table(title,headers,rows,ap,nota,kicker=None,widths=None):
    s=prs.slides.add_slide(BLANK); _header(s,title,PRES[ap],kicker)
    nr,nc=len(rows)+1,len(headers)
    gf=s.shapes.add_table(nr,nc,Inches(1.1),Inches(1.85),Inches(11.1),Inches(0.5))
    t=gf.table
    if widths:
        for j,w in enumerate(widths): t.columns[j].width=Inches(w)
    for j,h in enumerate(headers):
        c=t.cell(0,j); c.text=h; c.fill.solid(); c.fill.fore_color.rgb=MID
        p=c.text_frame.paragraphs[0]; p.font.bold=True; p.font.color.rgb=WHITE; p.font.size=Pt(17)
    for i,row in enumerate(rows,1):
        for j,val in enumerate(row):
            c=t.cell(i,j); c.text=val
            c.fill.solid(); c.fill.fore_color.rgb=(WHITE if i%2 else LGREY)
            p=c.text_frame.paragraphs[0]; p.font.size=Pt(15); p.font.color.rgb=INK
    META.append((s,ap)); notes(s,nota); return s

def closing():
    s=prs.slides.add_slide(BLANK); _rect(s,0,0,SW,SH,DARK)
    _rect(s,0,0,Inches(0.28),SH,ORANGE)
    tf=_tb(s,Inches(0.95),Inches(2.7),Inches(11.5),Inches(1.5))
    _para(tf,"Obrigado!",52,WHITE,bold=True,first=True,after=6)
    _para(tf,"Perguntas?",26,RGBColor(0xC7,0xDD,0xF0))
    META.append((s,None)); notes(s,"Agradecam e abram para perguntas. Tenham a "
        "cola e o glossario (docs/) a mao para duvidas tecnicas.")
    return s

def finalize():
    total=len(META)
    for i,(s,ap) in enumerate(META):
        if ap is None: continue
        tf=_tb(s,Inches(0.6),Inches(7.02),Inches(8),Inches(0.35))
        _para(tf,f"SAP-1  ·  {DISCIPLINA}"+("" if ap=="Todos" else f"  ·  {ap}"),
              10.5,GREY,first=True,after=0)
        tf=_tb(s,Inches(11.2),Inches(7.02),Inches(1.55),Inches(0.35))
        _para(tf,f"{i+1:02d} / {total:02d}",10.5,GREY,align=PP_ALIGN.RIGHT,first=True,after=0)

# =============================================================
# CONTEUDO
# =============================================================
cover("SAP-1 — Simple-As-Possible",
      "Arquitetura e implementação de um processador de 8 bits em Verilog · FPGA DE10-Lite")
notes(prs.slides[0],"Abertura. Apresentem-se, digam a disciplina e que vao mostrar a "
      "arquitetura completa do processador SAP-1 e sua implementacao, em 3 partes.")

agenda([
    (BLUE,"Introdução e Arquitetura","Pessoa 1",
     ["O que é o SAP-1","Objetivos","Especificações","Datapath e barramento"]),
    (TEAL,"Componentes e Implementação","Pessoa 2",
     ["PC, MAR, RAM, IR","Acumulador, B e ULA","Padrão de registrador","Saída e controle"]),
    (ORANGE,"Funcionamento e Resultados","Pessoa 3",
     ["Conjunto de instruções","Ciclo e máquina de estados","Palavra de controle",
      "Verificação e placa"]),
])
notes(prs.slides[1],"Apresentem rapidamente o roteiro: 3 partes. Cada um diz o que vai "
      "cobrir. Isso da uma visao geral antes de entrar nos detalhes.")

# ===================== PARTE 1 =====================
section(1,"Introdução e Arquitetura","Pessoa 1",
    ["O que é o SAP-1 e por que estudá-lo","Objetivos do projeto",
     "Especificações técnicas","O datapath e o barramento W"])

content("O que é o SAP-1",
    ["Processador didático de 8 bits, proposto por Albert Malvino",
     "'Simple-As-Possible': o menor computador que ainda executa programas",
     "Arquitetura de Von Neumann — programa e dados na mesma memória",
     "Base para entender processadores reais (busca, decodificação, execução)"],
    "Pessoa 1",
    "Contextualizem: o SAP-1 e o processador mais simples que ainda e um computador "
    "de verdade. Apesar de simples, ele tem os mesmos principios de um processador "
    "real: busca, decodifica e executa instrucoes.",
    kicker="Introdução",
    takeaway="Simples o bastante para entender por inteiro; completo o bastante para ensinar como um computador funciona.")

content("Objetivos do projeto",
    ["Compreender a organização interna de um processador (datapath + controle)",
     "Descrever o SAP-1 em Verilog RTL, módulo a módulo",
     "Verificar o funcionamento por simulação (testbenches)",
     "Sintetizar e executar na FPGA DE10-Lite, com saída visível nos displays"],
    "Pessoa 1",
    "Nossos objetivos foram quatro: entender a arquitetura, implementar em Verilog, "
    "verificar por simulacao e, por fim, rodar na placa fisica. Vamos seguir mais ou "
    "menos essa ordem na apresentacao.",
    kicker="Introdução")

table("Especificações técnicas",
    ["Característica","Valor"],
    [["Largura de dados","8 bits"],
     ["Memória","16 palavras × 8 bits (RAM 16×8)"],
     ["Endereçamento","4 bits (direto)"],
     ["Registradores de dados","Acumulador A e Registrador B"],
     ["ULA","Soma e subtração (complemento de dois)"],
     ["Instruções","5 — LDA, ADD, SUB, OUT, HLT"],
     ["Comunicação interna","Barramento único de 8 bits (W)"]],
    "Pessoa 1",
    "Aqui as especificacoes em numeros. O ponto principal: tudo e de 8 bits, memoria "
    "pequena de 16 posicoes, dois registradores de dados e apenas 5 instrucoes. Tudo "
    "conversa por um unico barramento.",
    kicker="Arquitetura", widths=[5.3,5.8])

image("O datapath — diagrama de blocos","img/datapath.png",
    "Blocos ligados ao barramento W. A Unidade de Controle gera os sinais (enables) que coordenam tudo.",
    "Pessoa 1",
    "Este e o mapa da arquitetura. No centro, o barramento. Em cima, busca e memoria: "
    "PC, MAR, RAM, IR. Embaixo, aritmetica: acumulador, ULA, registrador B, e a saida. "
    "A direita, a unidade de controle. As setas vermelhas sao os sinais de controle.",
    kicker="Arquitetura")

content("O barramento W e os 'enables'",
    ["Todos os blocos compartilham um único fio de 8 bits",
     "Sinais de 'falar' (colocam valor no barramento):",
     ("Ep (PC) · ~CE (RAM) · ~Ei (operando) · Ea (A) · Eu (ULA)",1),
     "Sinais de 'ouvir' (copiam do barramento):",
     ("~Lm (MAR) · ~Li (IR) · ~La (A) · ~Lb (B) · ~Lo (saída)",1),
     "Implementado como MULTIPLEXADOR (evita conflito elétrico)"],
    "Pessoa 1",
    "O barramento e a espinha dorsal. Como e um fio so, so um bloco pode escrever por "
    "vez. Os sinais de controle sao justamente os 'enables' de falar e ouvir. Passo "
    "para a Pessoa 2 detalhar cada bloco.",
    kicker="Arquitetura",
    takeaway="Regra de ouro — um bloco fala, um ou mais ouvem, por estado.")

# ===================== PARTE 2 =====================
section(2,"Componentes e Implementação","Pessoa 2",
    ["Os registradores: PC, MAR, IR","A memória RAM",
     "Acumulador, Registrador B e ULA","O padrão de projeto em Verilog"])

cards("Registradores de controle de fluxo",
    [("PC — Program Counter","Contador de 4 bits com o endereço da PRÓXIMA instrução. "
      "Cp incrementa; Ep coloca no barramento."),
     ("MAR — Memory Address Register","Guarda o endereço que a RAM vai ler. "
      "Carregado do barramento por ~Lm."),
     ("IR — Instruction Register","Guarda a instrução e a separa: opcode (o quê) "
      "e operando (onde). Carregado por ~Li."),
     ("Papel conjunto","PC diz 'qual instrução'; MAR 'aponta' na memória; "
      "IR 'segura' a instrução em execução.")],
    "Pessoa 2",
    "Estes tres registradores cuidam do fluxo. O PC aponta pra proxima instrucao, o MAR "
    "aponta pra memoria, e o IR guarda a instrucao atual, separando em opcode e "
    "operando. Juntos, eles fazem o processador saber o que executar e de onde.",
    kicker="Componentes")

content("RAM 16×8 — a memória",
    ["16 posições de 8 bits; guarda programa E dados (Von Neumann)",
     "Somente leitura durante a execução (o SAP-1 básico não tem STA)",
     "Formato de cada palavra:  [ opcode (4 bits) | operando (4 bits) ]",
     "O operando é um ENDEREÇO — endereçamento direto",
     ("Ex.: 'LDA 11' → carrega o conteúdo do endereço 11, não o número 11",1)],
    "Pessoa 2",
    "A memoria guarda 16 palavras. Ponto que confunde: o operando de uma instrucao e um "
    "endereco, nao um valor. 'LDA 11' carrega o que esta guardado no endereco 11. Isso "
    "e enderecamento direto.",
    kicker="Componentes",
    takeaway="Programa nas primeiras posições, dados nas últimas — mas a RAM não distingue: quem decide é a ordem de execução.")

cards("Núcleo aritmético",
    [("Acumulador A","Registrador principal — é onde a conta acontece. Alimenta a ULA "
      "e recebe os resultados. Sai no barramento por Ea."),
     ("Registrador B","Guarda o segundo operando da soma/subtração. Carregado por ~Lb."),
     ("ULA (somador/subtrator)","Combinacional: calcula A ± B. Su=0 soma, Su=1 subtrai "
      "(complemento de dois)."),
     ("Importante","A ULA não guarda nada — o resultado só vai ao barramento quando "
      "Eu está ativo.")],
    "Pessoa 2",
    "Aqui e o coracao aritmetico. O acumulador A guarda as contas, o B guarda o segundo "
    "numero, e a ULA soma ou subtrai. A ULA e combinacional: nao armazena, so mostra o "
    "resultado, que vai pro barramento quando o Eu liga.",
    kicker="Componentes")

code_slide("Padrão de projeto: o registrador",
    "Cinco blocos (MAR, IR, A, B e Saída) são o MESMO circuito — muda só a largura e o sinal de carga:",
    "always @(posedge clk or negedge clr_bar) begin\n"
    "    if (!clr_bar)        q <= 0;        // reset assíncrono\n"
    "    else if (!load_n)    q <= bus_in;   // carrega do barramento\n"
    "    // senão: mantém o valor guardado\n"
    "end",
    "Pessoa 2",
    "Um detalhe de implementacao que vale mostrar: cinco dos nossos blocos sao o mesmo "
    "circuito de registrador. Ele reseta de forma assincrona, carrega do barramento "
    "quando o sinal de carga esta ativo, e senao mantem o valor. Muda so a largura e o "
    "nome do sinal. Isso deixou o codigo limpo e reutilizavel.",
    takeaway="Reuso de um mesmo padrão RTL → código mais limpo e menos propenso a erro.",
    kicker="Implementação")

# ===================== PARTE 3 =====================
section(3,"Funcionamento e Resultados","Pessoa 3",
    ["Conjunto de instruções e ciclo de execução","A máquina de estados e a palavra de controle",
     "Exemplo de execução passo a passo","Verificação, placa e resultados"])

table("Conjunto de instruções",
    ["Instrução","Opcode","Operação","Estados úteis"],
    [["LDA n","0000","A ← memória[n]","T1–T5"],
     ["ADD n","0001","A ← A + memória[n]","T1–T6"],
     ["SUB n","0010","A ← A − memória[n]","T1–T6"],
     ["OUT","1110","saída ← A","T1–T4"],
     ["HLT","1111","para o processador","T1–T4"]],
    "Pessoa 3",
    "Sao 5 instrucoes. LDA carrega da memoria, ADD e SUB fazem aritmetica, OUT mostra o "
    "resultado e HLT para. Note que cada uma usa um numero diferente de estados uteis, "
    "mas todas ocupam os 6 estados de tempo.",
    kicker="Funcionamento", widths=[2.2,2.0,4.6,2.3])

content("O ciclo de instrução — 6 estados",
    ["Toda instrução é executada em 6 estados de tempo (T1 a T6)",
     "Ciclo de BUSCA (T1–T3) — igual para toda instrução:",
     ("T1: MAR ← PC    ·    T2: PC ← PC+1    ·    T3: IR ← RAM",1),
     "Ciclo de EXECUÇÃO (T4–T6) — depende do opcode",
     "Temporização: estados avançam na descida do clock; registradores carregam na subida"],
    "Pessoa 3",
    "Cada instrucao leva 6 passos. Os tres primeiros sao a busca, sempre iguais. Os tres "
    "ultimos sao a execucao, que muda conforme a instrucao. Ha um cuidado de "
    "temporizacao: os estados avancam na descida do clock, e os registradores carregam "
    "na subida, pra tudo ficar estavel.",
    kicker="Funcionamento",
    takeaway="Busca (T1–T3) sempre igual; execução (T4–T6) específica de cada instrução.")

image("A máquina de estados (FSM)","../fsm/fsm_diagram.png",
    "Contador de anel IDLE → T1 → … → T6 → T1. O reset entra em IDLE; o HLT congela em T4.",
    "Pessoa 3",
    "A unidade de controle e, na pratica, esta maquina de estados: um contador de anel "
    "que percorre T1 a T6 e volta. O reset entra no IDLE e o HLT congela no T4. Cada "
    "estado corresponde a uma acao no datapath.",
    kicker="Controle")

content("A palavra de controle (12 bits)",
    ["A cada estado, o controlador gera 12 sinais de uma vez (a palavra CON)",
     "Cada bit liga/desliga um bloco (falar ou ouvir no barramento)",
     ("Ativos-alto (1 = liga):  Cp  Ep  Ea  Su  Eu",1),
     ("Ativos-baixo (0 = liga):  ~Lm  ~CE  ~Li  ~Ei  ~La  ~Lb  ~Lo",1),
     "Depende apenas do estado e do opcode — nunca do dado"],
    "Pessoa 3",
    "O que sai da maquina de estados e a palavra de controle: 12 bits que dizem, a cada "
    "estado, o que cada bloco faz. Alguns sinais sao ativos em baixo. E o mais "
    "importante: essa palavra depende so do estado e da instrucao, nunca do dado. Por "
    "isso o comportamento e sempre previsivel, deterministico.",
    kicker="Controle",
    takeaway="A palavra de controle é a 'receita' de cada estado — determinística.")

image("Microoperações por instrução","../fsm/fsm_exec_diagram.png",
    "Busca (azul) igual para todas; execução (verde) específica. Cada coluna é uma instrução.",
    "Pessoa 3",
    "Este diagrama mostra a fase de execucao de cada instrucao. A busca, em azul, e "
    "sempre igual. A execucao, em verde, muda: LDA carrega o acumulador, ADD e SUB usam "
    "a ULA, OUT manda pra saida e HLT para.",
    kicker="Controle")

table("Exemplo de execução: ADD",
    ["Estado","Sinais ativos","Transferência no barramento"],
    [["T1","Ep, ~Lm","MAR ← PC"],
     ["T2","Cp","PC ← PC + 1"],
     ["T3","~CE, ~Li","IR ← RAM (a instrução ADD)"],
     ["T4","~Ei, ~Lm","MAR ← operando (endereço do dado)"],
     ["T5","~CE, ~Lb","B ← RAM (o dado)"],
     ["T6","Eu, ~La","A ← A + B"]],
    "Pessoa 3",
    "Vamos juntar tudo num exemplo real: o ADD. Cada linha e um estado, com os sinais "
    "que ligam e a transferencia que acontece no barramento. Nos tres primeiros ele "
    "busca; no T4 acha o dado; no T5 carrega o B; e no T6 soma. Seis passos e a soma "
    "esta feita.",
    kicker="Funcionamento", widths=[1.8,3.3,6.0])

cards("Implementação em Verilog",
    [("11 módulos","Um módulo por bloco (PC, MAR, RAM, IR, A, B, ULA, saída, controle…) "
      "mais o topo que conecta tudo pelo barramento."),
     ("Barramento por mux","Em vez de tri-state interno — mais robusto e sintetizável "
      "em FPGA."),
     ("HLT por clock-enable","Congela o contador de anel sem desligar o clock (evita "
      "gating de clock)."),
     ("Reset seguro","Assíncrono no assert, sincronizado no desassert (sem violação de "
      "recovery/removal).")],
    "Pessoa 3",
    "Na implementacao, cada bloco virou um modulo Verilog. Tomamos cuidado com boas "
    "praticas de projeto digital: barramento por multiplexador, HLT por clock-enable em "
    "vez de desligar o clock, e um reset bem tratado. Isso deixa o projeto robusto.",
    kicker="Implementação")

content("Verificação por simulação",
    ["Simulação no ModelSim, com testbenches próprios",
     "14 testbenches autoverificáveis:",
     ("um por componente (PC, MAR, ULA, controlador, …)",1),
     ("+ um de integração, que roda um programa e confere o resultado final",1),
     "Resultado: todos passam (PASSOU · 0 erros)"],
    "Pessoa 3",
    "Nao confiamos so no olho: cada bloco tem um teste automatico que diz PASSOU ou "
    "FALHOU. Sao 14 testes no total, incluindo um que roda um programa inteiro e confere "
    "o resultado. Todos passam - isso da seguranca de que o processador esta correto.",
    kicker="Resultados",
    takeaway="Verificação automatizada — 14/14 testbenches aprovados.")

content("Execução na placa — DE10-Lite",
    ["Roda de verdade na FPGA Intel MAX 10, com clock lento (~1 Hz)",
     "Displays de 7 segmentos revelam o funcionamento interno:",
     ("HEX5 = estado (1–6)   ·   HEX4 = instrução (L/A/5/o/H)",1),
     ("HEX3-2 = acumulador A   ·   HEX1-0 = resultado",1),
     "Modo passo a passo por botão — avança um estado de cada vez"],
    "Pessoa 3",
    "E roda na placa fisica. O clock e bem lento pra dar pra ver acontecendo. Os seis "
    "displays mostram o estado, a instrucao, o acumulador e o resultado. Da pra avancar "
    "passo a passo com um botao. Se possivel, facam a demonstracao ao vivo agora.",
    kicker="Resultados")

table("Programas de teste executados",
    ["Programa","Técnica","Resultado"],
    [["3³","Potência por somas repetidas","27  (0x1B)"],
     ["(((7+3)−2)+5)−4 +7−3+5","Soma/subtração encadeadas","18  (0x12)"],
     ["3 × 4","Multiplicação por somas","12  (0x0C)"],
     ["12 ÷ 4","Divisão por subtrações","3  (0x03)"]],
    "Pessoa 3",
    "Testamos com quatro programas. Como o SAP-1 so soma e subtrai, fazemos "
    "multiplicacao com somas repetidas e divisao com subtracoes. Todos deram o resultado "
    "esperado, tanto na simulacao quanto na placa.",
    kicker="Resultados", widths=[4.6,4.3,2.4])

content("Conclusão",
    ["Implementamos um processador completo, do zero, em Verilog RTL",
     "Cobrimos a arquitetura (datapath, barramento, controle) e o funcionamento "
     "(busca–execução, palavra de controle)",
     "Aplicamos boas práticas de projeto digital",
     "Verificado por simulação e validado na placa real",
     ("O SAP-1 mostra, de forma tangível, como um computador realmente funciona",1)],
    "Pessoa 3",
    "Para concluir: construimos um processador inteiro e entendemos tanto a arquitetura "
    "quanto o funcionamento. Usamos boas praticas, verificamos com testes automaticos e "
    "validamos na placa. Foi um projeto que uniu teoria e pratica.",
    kicker="Encerramento")

content("Trabalhos futuros",
    ["Instrução de escrita na memória (STA) → programas que alteram dados",
     "Instruções de desvio (JMP / JZ) → laços e decisões reais",
     "Ampliar memória e conjunto de instruções (rumo ao SAP-2)",
     "Otimização de temporização e relatório de Fmax (análise de timing)"],
    "Pessoa 3",
    "Como evolucao, o SAP-1 pode ganhar uma instrucao de escrita, desvios pra fazer "
    "lacos de verdade, e mais instrucoes - o proximo passo seria o SAP-2. Tambem daria "
    "pra fazer a analise de timing completa.",
    kicker="Encerramento")

content("Referências",
    ["MALVINO, A. P.; BROWN, J. A. Eletrônica Digital / Digital Computer Electronics. "
     "3ª ed. — capítulo do SAP-1",
     "INTEL. Quartus Prime Lite 20.1 — documentação e guias de uso",
     "TERASIC. DE10-Lite User Manual (FPGA Intel MAX 10 10M50DAF484C7G)",
     "IEEE Std 1364 — linguagem Verilog HDL"],
    "Pessoa 3",
    "Nossas referencias principais: o livro do Malvino, onde o SAP-1 e definido; a "
    "documentacao do Quartus e da placa DE10-Lite; e o padrao da linguagem Verilog.",
    kicker="Encerramento")

closing()
finalize()

prs.save("SAP1_apresentacao.pptx")
print("PPTX gerado: SAP1_apresentacao.pptx (%d slides)"%len(prs.slides._sldIdLst))
