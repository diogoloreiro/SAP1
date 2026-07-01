#!/usr/bin/env python3
# =============================================================
# gerar_pptx_4partes.py - Apresentacao do SAP-1 em 4 PARTES,
# apresentada por 3 PESSOAS (a Pessoa 2 cobre as partes 2 e 3):
#   1) Objetivo e FSM ............ Pessoa 1
#   2) PC e Memoria RAM .......... Pessoa 2
#   3) Palavra de controle ....... Pessoa 2
#   4) Validacao e conclusao ..... Pessoa 3
# Uso (na pasta apresentacao):  python3 gerar_pptx_4partes.py
# =============================================================
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import math

DISCIPLINA = "Arquitetura de Computadores"      # <-- ajuste
INSTITUICAO = "[Instituição] · [Curso]"          # <-- ajuste
DATA = "2025"                                     # <-- ajuste

DARK  = RGBColor(0x0D,0x2B,0x4E); MID = RGBColor(0x0D,0x3B,0x66)
BLUE  = RGBColor(0x15,0x65,0xC0); TEAL= RGBColor(0x00,0x83,0x8F)
GREY  = RGBColor(0x5A,0x63,0x6E); LGREY=RGBColor(0xEC,0xF1,0xF5)
WHITE = RGBColor(0xFF,0xFF,0xFF); INK = RGBColor(0x1B,0x26,0x31)
ORANGE=RGBColor(0xE6,0x5C,0x00)
CODEBG= RGBColor(0x1E,0x24,0x30); CODEFG=RGBColor(0xE6,0xED,0xF3)
PRES  = {"Pessoa 1": BLUE, "Pessoa 2": TEAL, "Pessoa 3": ORANGE}
NPARTS = 4

prs = Presentation()
prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height
META = []

def _tb(s,l,t,w,h):
    tf=s.shapes.add_textbox(l,t,w,h).text_frame; tf.word_wrap=True
    tf.margin_left=Pt(2); tf.margin_right=Pt(2); return tf
def _rect(s,l,t,w,h,c,rounded=False):
    shp=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,l,t,w,h)
    shp.fill.solid(); shp.fill.fore_color.rgb=c; shp.line.fill.background()
    shp.shadow.inherit=False; return shp
def _para(tf,text,size,color,bold=False,italic=False,align=None,after=6,first=False,level=0,name=None):
    p=tf.paragraphs[0] if first else tf.add_paragraph()
    p.text=text; p.font.size=Pt(size); p.font.color.rgb=color
    p.font.bold=bold; p.font.italic=italic; p.level=level
    if name: p.font.name=name
    if align is not None: p.alignment=align
    p.space_after=Pt(after); return p
def notes(s,t): s.notes_slide.notes_text_frame.text=t

def _header(s,title,cor,kicker=None):
    _rect(s,0,0,SW,Inches(1.2),MID); _rect(s,0,Inches(1.2),SW,Inches(0.07),cor)
    if kicker:
        _para(_tb(s,Inches(0.62),Inches(0.16),Inches(11),Inches(0.35)),
              kicker.upper(),12,RGBColor(0x8F,0xB8,0xDE),bold=True,after=0,first=True)
        _para(_tb(s,Inches(0.6),Inches(0.5),Inches(12.1),Inches(0.62)),
              title,26,WHITE,bold=True,after=0,first=True)
    else:
        tf=_tb(s,Inches(0.6),Inches(0.2),Inches(12.1),Inches(0.8)); tf.vertical_anchor=MSO_ANCHOR.MIDDLE
        _para(tf,title,28,WHITE,bold=True,after=0,first=True)

def cover(title,sub):
    s=prs.slides.add_slide(BLANK); _rect(s,0,0,SW,SH,DARK); _rect(s,0,0,Inches(0.28),SH,ORANGE)
    _para(_tb(s,Inches(0.95),Inches(0.85),Inches(11.4),Inches(0.4)),
          DISCIPLINA.upper(),14,RGBColor(0x8F,0xB8,0xDE),bold=True,first=True)
    tf=_tb(s,Inches(0.9),Inches(1.7),Inches(11.6),Inches(1.9))
    _para(tf,title,48,WHITE,bold=True,first=True,after=4)
    _para(tf,sub,20,RGBColor(0xC7,0xDD,0xF0))
    _rect(s,Inches(0.95),Inches(4.35),Inches(4.2),Inches(0.045),ORANGE)
    tf=_tb(s,Inches(0.92),Inches(4.6),Inches(11.5),Inches(1.7))
    for i in range(3):
        _para(tf,f"Integrante {i+1}  ·  [nome]",18,WHITE,first=(i==0),after=5)
    _para(_tb(s,Inches(0.92),Inches(6.7),Inches(11.5),Inches(0.5)),
          f"{INSTITUICAO}   ·   {DATA}",14,RGBColor(0x9F,0xB6,0xCC),first=True)
    META.append((s,None)); return s

def section(num,tit,ap,topics):
    s=prs.slides.add_slide(BLANK); cor=PRES[ap]; _rect(s,0,0,SW,SH,cor)
    _para(_tb(s,Inches(0.95),Inches(1.9),Inches(11),Inches(0.5)),
          f"PARTE {num} DE {NPARTS}",16,WHITE,bold=True,first=True)
    _para(_tb(s,Inches(0.9),Inches(2.5),Inches(11.4),Inches(1.3)),
          tit,40,WHITE,bold=True,first=True)
    _rect(s,Inches(0.95),Inches(3.95),Inches(3.6),Inches(0.05),WHITE)
    tf=_tb(s,Inches(0.95),Inches(4.2),Inches(11),Inches(2.2))
    for i,t in enumerate(topics): _para(tf,"•  "+t,18,WHITE,first=(i==0),after=7)
    _para(_tb(s,Inches(0.95),Inches(6.5),Inches(11),Inches(0.5)),
          "Apresenta: "+ap,17,RGBColor(0xEE,0xEE,0xEE),bold=True,first=True)
    META.append((s,None)); return s

def agenda(parts):
    s=prs.slides.add_slide(BLANK); _header(s,"Sumário",BLUE)
    n=len(parts); marg=Inches(0.6); gap=Inches(0.22)
    cw=(SW-2*marg-gap*(n-1))/n; top=Inches(1.75); ch=Inches(4.95)
    for i,(cor,titulo,ap,itens) in enumerate(parts):
        lx=marg+(cw+gap)*i
        _rect(s,lx,top,cw,ch,LGREY,rounded=True)
        _rect(s,lx,top,cw,Inches(1.0),cor,rounded=True)
        tf=_tb(s,lx+Inches(0.2),top+Inches(0.12),cw-Inches(0.4),Inches(0.85))
        _para(tf,f"PARTE {i+1}",11,WHITE,bold=True,first=True,after=0)
        _para(tf,titulo,14.5,WHITE,bold=True,after=0)
        tf=_tb(s,lx+Inches(0.22),top+Inches(1.2),cw-Inches(0.4),Inches(3.1))
        for j,it in enumerate(itens): _para(tf,"•  "+it,12,INK,first=(j==0),after=6)
        _para(_tb(s,lx+Inches(0.22),top+Inches(4.4),cw-Inches(0.4),Inches(0.45)),
              ap,12.5,cor,bold=True,first=True)
    META.append((s,"Todos")); return s

def content(title,bullets,ap,nota,kicker=None,takeaway=None):
    s=prs.slides.add_slide(BLANK); _header(s,title,PRES[ap],kicker)
    bh=Inches(4.6) if takeaway else Inches(5.1)
    body=_tb(s,Inches(0.75),Inches(1.55),Inches(11.8),bh)
    for i,it in enumerate(bullets):
        txt,lvl=(it if isinstance(it,tuple) else (it,0))
        _para(body,("▸  " if lvl==0 else "–  ")+txt,20 if lvl==0 else 17,
              INK if lvl==0 else GREY,first=(i==0),after=10,level=lvl)
    if takeaway:
        _rect(s,Inches(0.75),Inches(6.15),Inches(11.85),Inches(0.72),LGREY,rounded=True)
        _rect(s,Inches(0.75),Inches(6.15),Inches(0.12),Inches(0.72),PRES[ap])
        tf=_tb(s,Inches(1.05),Inches(6.2),Inches(11.4),Inches(0.62)); tf.vertical_anchor=MSO_ANCHOR.MIDDLE
        _para(tf,"Chave:  "+takeaway,15,MID,bold=True,first=True,after=0)
    META.append((s,ap)); notes(s,nota); return s

def cards(title,items,ap,nota,kicker=None):
    s=prs.slides.add_slide(BLANK); _header(s,title,PRES[ap],kicker)
    n=len(items); cols=2 if n<=4 else 3; rows=math.ceil(n/cols)
    x0=Inches(0.7); y0=Inches(1.65)
    cw=(SW-Inches(1.4)-Inches(0.3)*(cols-1))/cols
    ch=(Inches(5.15)-Inches(0.3)*(rows-1))/rows
    for i,(t,body) in enumerate(items):
        r,c=divmod(i,cols); lx=x0+(cw+Inches(0.3))*c; ty=y0+(ch+Inches(0.3))*r
        _rect(s,lx,ty,cw,ch,LGREY,rounded=True); _rect(s,lx,ty,cw,Inches(0.12),PRES[ap])
        tf=_tb(s,lx+Inches(0.22),ty+Inches(0.2),cw-Inches(0.44),ch-Inches(0.4))
        _para(tf,t,17,MID,bold=True,first=True,after=5); _para(tf,body,14,GREY,after=0)
    META.append((s,ap)); notes(s,nota); return s

def image(title,img,cap,ap,nota,kicker=None):
    s=prs.slides.add_slide(BLANK); _header(s,title,PRES[ap],kicker)
    pic=s.shapes.add_picture(img,0,Inches(1.5),height=Inches(4.95)); pic.left=int((SW-pic.width)/2)
    if cap: _para(_tb(s,Inches(0.6),Inches(6.55),Inches(12.1),Inches(0.4)),
                  cap,13,GREY,italic=True,align=PP_ALIGN.CENTER,first=True,after=0)
    META.append((s,ap)); notes(s,nota); return s

def table(title,headers,rows,ap,nota,kicker=None,widths=None):
    s=prs.slides.add_slide(BLANK); _header(s,title,PRES[ap],kicker)
    nr,nc=len(rows)+1,len(headers)
    t=s.shapes.add_table(nr,nc,Inches(1.1),Inches(1.8),Inches(11.1),Inches(0.5)).table
    if widths:
        for j,w in enumerate(widths): t.columns[j].width=Inches(w)
    for j,h in enumerate(headers):
        c=t.cell(0,j); c.text=h; c.fill.solid(); c.fill.fore_color.rgb=MID
        p=c.text_frame.paragraphs[0]; p.font.bold=True; p.font.color.rgb=WHITE; p.font.size=Pt(16)
    for i,row in enumerate(rows,1):
        for j,val in enumerate(row):
            c=t.cell(i,j); c.text=val; c.fill.solid(); c.fill.fore_color.rgb=(WHITE if i%2 else LGREY)
            p=c.text_frame.paragraphs[0]; p.font.size=Pt(14); p.font.color.rgb=INK
    META.append((s,ap)); notes(s,nota); return s

def closing():
    s=prs.slides.add_slide(BLANK); _rect(s,0,0,SW,SH,DARK); _rect(s,0,0,Inches(0.28),SH,ORANGE)
    tf=_tb(s,Inches(0.95),Inches(2.7),Inches(11.5),Inches(1.5))
    _para(tf,"Obrigado!",52,WHITE,bold=True,first=True,after=6)
    _para(tf,"Perguntas?",26,RGBColor(0xC7,0xDD,0xF0))
    META.append((s,None)); notes(s,"Agradecam e abram para perguntas."); return s

def finalize():
    total=len(META)
    for i,(s,ap) in enumerate(META):
        if ap is None: continue
        _para(_tb(s,Inches(0.6),Inches(7.02),Inches(8),Inches(0.35)),
              f"SAP-1  ·  {DISCIPLINA}"+("" if ap=="Todos" else f"  ·  {ap}"),
              10.5,GREY,first=True,after=0)
        _para(_tb(s,Inches(11.2),Inches(7.02),Inches(1.55),Inches(0.35)),
              f"{i+1:02d} / {total:02d}",10.5,GREY,align=PP_ALIGN.RIGHT,first=True,after=0)

# =============================================================
# CONTEUDO — 4 PARTES, 3 APRESENTADORES
# =============================================================
cover("SAP-1 — Simple-As-Possible",
      "Arquitetura e implementação de um processador de 8 bits em Verilog · FPGA DE10-Lite")
notes(prs.slides[0],"Abertura. Apresentem-se e digam que vao mostrar o processador SAP-1, "
      "em 4 partes, por 3 integrantes.")

agenda([
    (BLUE,  "Objetivo e FSM","Pessoa 1",
     ["Objetivo do trabalho","Nossa abordagem","Arquitetura (datapath)","Ciclo e máquina de estados"]),
    (TEAL,  "Program Counter e RAM","Pessoa 2",
     ["Program Counter (PC)","MAR — endereço","Memória RAM 16×8","Ciclo de busca"]),
    (TEAL,  "Palavra de controle","Pessoa 2",
     ["Os 12 sinais de controle","Falar e ouvir","Palavras por estado","Exemplo: ADD"]),
    (ORANGE,"Validação e conclusão","Pessoa 3",
     ["Implementação em Verilog","Verificação (testbenches)","Placa e programas","Conclusão"]),
])
notes(prs.slides[1],"Mostrem o roteiro: 4 partes. Pessoa 1 abre; Pessoa 2 cobre as partes "
      "2 e 3 (componentes e controle); Pessoa 3 fecha com validacao e conclusao.")

# ============ PARTE 1 — OBJETIVO E FSM (Pessoa 1) ============
section(1,"Objetivo e FSM","Pessoa 1",
    ["O objetivo do trabalho e nossa abordagem","A arquitetura que implementamos",
     "O ciclo de instrução (6 estados)","A máquina de estados (FSM)"])

content("Objetivo do trabalho",
    ["Implementar o processador SAP-1 (do Malvino) em Verilog",
     "Processador didático de 8 bits, com 5 instruções, arquitetura de Von Neumann",
     "Entender, na prática, como um computador busca e executa instruções",
     "Sintetizar e rodar na FPGA DE10-Lite, com saída visível nos displays"],
    "Pessoa 1",
    "Comecem pelo objetivo: implementar o SAP-1, o processador mais simples que ainda e "
    "um computador de verdade. A ideia era entender, construindo, como um processador "
    "funciona por dentro - e rodar de verdade na placa.",
    kicker="Introdução",
    takeaway="Construir um processador do zero para entender como um computador funciona.")

content("Nossa abordagem — como pensamos o trabalho",
    ["Dividir o processador em módulos: um bloco = um arquivo Verilog",
     "Controlar tudo com uma MÁQUINA DE ESTADOS (FSM):",
     ("a FSM percorre os passos de cada instrução e comanda os blocos",1),
     "Verificar cada parte por simulação antes de juntar (testbenches)",
     "Só então sintetizar e testar na placa"],
    "Pessoa 1",
    "Nossa estrategia teve quatro pilares: primeiro, modularizar - cada bloco virou um "
    "arquivo. Segundo, e o mais importante, controlar tudo por uma maquina de estados. "
    "Terceiro, verificar cada parte por simulacao. E so no fim ir pra placa. Essa ordem "
    "evitou erros.",
    kicker="Metodologia",
    takeaway="Modularizar + controlar por FSM + verificar em cada etapa.")

image("A arquitetura que implementamos","img/datapath.png",
    "Blocos ligados ao barramento W; a Unidade de Controle (a FSM) gera os sinais que coordenam tudo.",
    "Pessoa 1",
    "Esta e a arquitetura. No centro, o barramento por onde os dados passam. Em cima, "
    "busca e memoria. Embaixo, a aritmetica. E a direita, a unidade de controle - que e "
    "a nossa maquina de estados, o foco desta primeira parte.",
    kicker="Arquitetura")

content("O ciclo de instrução — 6 estados",
    ["Toda instrução é executada em 6 estados de tempo: T1 a T6",
     "Ciclo de BUSCA (T1–T3) — igual para toda instrução",
     ("T1: MAR ← PC   ·   T2: PC ← PC+1   ·   T3: IR ← RAM",1),
     "Ciclo de EXECUÇÃO (T4–T6) — depende da instrução",
     "É a FSM que percorre esses estados e decide o que fazer em cada um"],
    "Pessoa 1",
    "O funcionamento e organizado em 6 estados. Os tres primeiros sao a busca, sempre "
    "iguais: pegar a instrucao da memoria. Os tres ultimos sao a execucao. Quem controla "
    "essa sequencia e a maquina de estados.",
    kicker="Funcionamento",
    takeaway="6 estados por instrução: busca (T1–T3) + execução (T4–T6).")

image("A máquina de estados (FSM)","../fsm/fsm_diagram.png",
    "Contador de anel IDLE → T1 → … → T6 → T1. O reset entra em IDLE; o HLT congela em T4.",
    "Pessoa 1",
    "E aqui esta a FSM. E um contador de anel que gira pelos estados T1 a T6 e volta ao "
    "T1 para a proxima instrucao. O reset a coloca no IDLE, e a instrucao HLT a congela "
    "no T4. Cada estado corresponde a uma acao. Passo para a Pessoa 2 mostrar os "
    "primeiros blocos.",
    kicker="Controle")

# ============ PARTE 2 — PC E RAM (Pessoa 2) ============
section(2,"Program Counter e Memória RAM","Pessoa 2",
    ["O Program Counter (PC)","O MAR — registrador de endereço",
     "A memória RAM 16×8","Como eles trabalham juntos na busca"])

content("Program Counter (PC)",
    ["Contador de 4 bits que guarda o endereço da PRÓXIMA instrução",
     "Sinal Cp: incrementa o PC (+1) na borda de subida do clock",
     "Sinal Ep: coloca o valor do PC no barramento",
     ("No T2 de toda instrução, o Cp é ativado → o PC já aponta para a seguinte",1),
     "É de 4 bits porque a memória só tem 16 posições (0–15)"],
    "Pessoa 2",
    "O PC e um contador simples: guarda o endereco da proxima instrucao. O sinal Cp faz "
    "ele incrementar, e o Ep coloca o valor no barramento. Ele e de 4 bits porque a "
    "memoria tem so 16 posicoes.",
    kicker="Componente",
    takeaway="O PC é quem diz 'qual é a próxima instrução'.")

content("MAR — Registrador de Endereço de Memória",
    ["Guarda o endereço de 4 bits que a RAM vai ler",
     "Carregado do barramento quando ~Lm está ativo",
     "É o 'elo' entre o processador e a memória:",
     ("T1: recebe o PC (para buscar a instrução)",1),
     ("T4: recebe o operando do IR (para buscar o dado)",1)],
    "Pessoa 2",
    "O MAR e o elo com a memoria: ele segura o endereco que a RAM vai ler. Ele e "
    "carregado duas vezes por instrucao: no T1 recebe o PC, para buscar a instrucao; e "
    "no T4 recebe o operando, para buscar o dado.",
    kicker="Componente")

content("Memória RAM 16×8",
    ["16 posições de 8 bits; guarda PROGRAMA e DADOS juntos (Von Neumann)",
     "Somente leitura durante a execução (não há instrução de escrita, STA)",
     "Cada palavra:  [ opcode (4 bits) | operando (4 bits) ]",
     "O operando é um ENDEREÇO, não um valor → endereçamento direto",
     ("Ex.: 'LDA 11' carrega o conteúdo do endereço 11, não o número 11",1)],
    "Pessoa 2",
    "A memoria guarda 16 palavras de 8 bits, programa e dados juntos. Ponto que "
    "confunde: o operando de uma instrucao e um endereco, nao um valor. 'LDA 11' carrega "
    "o que esta no endereco 11. Isso e enderecamento direto.",
    kicker="Componente",
    takeaway="O número na instrução é 'onde está', não 'o quê' (endereçamento direto).")

table("Como PC, MAR e RAM trabalham na busca (T1–T3)",
    ["Estado","Ação","Sinais ativos"],
    [["T1","MAR ← PC (endereço da instrução)","Ep, ~Lm"],
     ["T2","PC ← PC + 1 (aponta para a próxima)","Cp"],
     ["T3","IR ← RAM (lê a instrução da memória)","~CE, ~Li"]],
    "Pessoa 2",
    "Juntando os tres: na busca, o PC manda o endereco pro MAR (T1), o PC se incrementa "
    "(T2), e a RAM entrega a instrucao, que vai pro IR (T3). Esses tres passos sao iguais "
    "para toda instrucao. Sigo agora para a palavra de controle, que gera esses sinais.",
    kicker="Funcionamento", widths=[1.8,6.6,2.7])

# ============ PARTE 3 — PALAVRA DE CONTROLE (Pessoa 2) ============
section(3,"Palavra de Controle","Pessoa 2",
    ["Os 12 sinais de controle","Falar e ouvir no barramento",
     "As palavras de controle por estado","Exemplo: execução do ADD"])

content("A palavra de controle (12 bits)",
    ["A cada estado, a FSM gera 12 sinais de uma vez — a palavra CON",
     "Cada bit liga/desliga um bloco (falar ou ouvir no barramento)",
     ("CON = { Cp, Ep, ~Lm, ~CE, ~Li, ~Ei, ~La, Ea, Su, Eu, ~Lb, ~Lo }",1),
     "Depende APENAS do estado e do opcode — nunca do dado",
     ("Por isso o processador é determinístico",1)],
    "Pessoa 2",
    "O que sai da maquina de estados e a palavra de controle: 12 bits que, a cada "
    "estado, dizem o que cada bloco faz. Cada bit comanda um bloco. E o mais importante: "
    "ela depende so do estado e da instrucao, nunca do dado - por isso o comportamento e "
    "sempre previsivel.",
    kicker="Controle",
    takeaway="A palavra de controle é a 'receita' de cada estado.")

content("Falar e ouvir no barramento",
    ["Cada estado escolhe UM bloco que fala e UM ou mais que ouvem",
     ("Falam: Ep (PC) · ~CE (RAM) · ~Ei (operando) · Ea (A) · Eu (ULA)",1),
     ("Ouvem: ~Lm (MAR) · ~Li (IR) · ~La (A) · ~Lb (B) · ~Lo (saída)",1),
     "Sinais ativos-ALTO (1=liga): Cp Ep Ea Su Eu",
     "Sinais ativos-BAIXO (0=liga): os com til (~)",
     ("Repouso (IDLE) = 0011_1110_0011 — nada ativo",1)],
    "Pessoa 2",
    "Na pratica, a palavra escolhe quem fala e quem ouve no barramento. Ha cinco sinais "
    "de falar e cinco de ouvir. Um detalhe: alguns sao ativos em baixo, agem quando "
    "valem zero. Por isso o estado de repouso nao e tudo zero.",
    kicker="Controle",
    takeaway="Um fala, um ou mais ouvem — a cada estado.")

table("As palavras de controle por estado",
    ["Estado","CON (12 bits)","Ação"],
    [["T1","0101_1110_0011","MAR ← PC"],
     ["T2","1011_1110_0011","PC ← PC+1"],
     ["T3","0010_0110_0011","IR ← RAM"],
     ["T4 (LDA/ADD/SUB)","0001_1010_0011","MAR ← operando"],
     ["T5 (LDA)","0010_1100_0011","A ← RAM"],
     ["T5 (ADD/SUB)","0010_1110_0001","B ← RAM"],
     ["T6 (ADD)","0011_1100_0111","A ← A+B"],
     ["T6 (SUB)","0011_1100_1111","A ← A−B"]],
    "Pessoa 2",
    "Esta tabela mostra as 'receitas' - a palavra de 12 bits de cada estado e a acao que "
    "ela produz. A busca (T1 a T3) e igual pra todas; a execucao (T4 a T6) muda conforme "
    "a instrucao. Da pra ler cada palavra comparando com o repouso.",
    kicker="Controle", widths=[3.0,4.2,3.9])

table("Exemplo: execução do ADD",
    ["Estado","Sinais ativos","Transferência"],
    [["T1","Ep, ~Lm","MAR ← PC"],
     ["T2","Cp","PC ← PC+1"],
     ["T3","~CE, ~Li","IR ← RAM (instrução ADD)"],
     ["T4","~Ei, ~Lm","MAR ← operando"],
     ["T5","~CE, ~Lb","B ← RAM (o dado)"],
     ["T6","Eu, ~La","A ← A + B"]],
    "Pessoa 2",
    "Juntando tudo num exemplo: o ADD, estado por estado. Cada linha e uma palavra de "
    "controle em acao. Busca nos tres primeiros, acha o dado no T4, carrega o B no T5, e "
    "soma no T6. Passo para a Pessoa 3 mostrar como validamos tudo isso.",
    kicker="Funcionamento", widths=[1.8,3.4,5.9])

# ============ PARTE 4 — VALIDACAO E CONCLUSAO (Pessoa 3) ============
section(4,"Validação e Conclusão","Pessoa 3",
    ["Implementação em Verilog","Verificação por simulação",
     "Execução na placa e programas testados","Conclusão e trabalhos futuros"])

cards("Implementação em Verilog",
    [("11 módulos","Um módulo por bloco + o topo que conecta tudo pelo barramento W."),
     ("Barramento por mux","Em vez de tri-state interno — robusto e sintetizável em FPGA."),
     ("HLT por clock-enable","Congela a FSM sem desligar o clock (evita gating de clock)."),
     ("Reset seguro","Assíncrono no assert, sincronizado no desassert.")],
    "Pessoa 3",
    "Na implementacao, cada bloco virou um modulo Verilog. Tomamos cuidado com boas "
    "praticas: barramento por multiplexador, HLT por clock-enable em vez de desligar o "
    "clock, e um reset bem tratado. Isso deixou o projeto robusto e sintetizavel.",
    kicker="Implementação")

content("Verificação por simulação",
    ["Simulação no ModelSim, com testbenches próprios",
     "14 testbenches autoverificáveis:",
     ("um por componente (PC, MAR, ULA, controlador, …)",1),
     ("+ um de integração, que roda um programa e confere o resultado final",1),
     "Resultado: todos passam (PASSOU · 0 erros)"],
    "Pessoa 3",
    "Para validar, nao confiamos so no olho: cada bloco tem um teste automatico que diz "
    "PASSOU ou FALHOU. Sao 14 testes, incluindo um que roda um programa inteiro e "
    "confere o resultado. Todos passam.",
    kicker="Validação",
    takeaway="Verificação automatizada — 14/14 testbenches aprovados.")

content("Execução na placa — DE10-Lite",
    ["Roda na FPGA Intel MAX 10, com clock lento (~1 Hz) para dar para ver",
     "Displays de 7 segmentos mostram o funcionamento interno:",
     ("HEX5 = estado (1–6) · HEX4 = instrução · HEX3-2 = A · HEX1-0 = resultado",1),
     "Modo passo a passo por botão (avança um estado por vez)"],
    "Pessoa 3",
    "E roda de verdade na placa, com clock lento pra dar pra acompanhar. Os seis displays "
    "mostram o estado, a instrucao, o acumulador e o resultado. Da pra avancar passo a "
    "passo. Se possivel, facam a demonstracao ao vivo.",
    kicker="Validação")

table("Programas testados",
    ["Programa","Técnica","Resultado"],
    [["3³","Potência por somas repetidas","27  (0x1B)"],
     ["(((7+3)−2)+5)−4 +7−3+5","Soma/subtração encadeadas","18  (0x12)"],
     ["3 × 4","Multiplicação por somas","12  (0x0C)"],
     ["12 ÷ 4","Divisão por subtrações","3  (0x03)"]],
    "Pessoa 3",
    "Testamos com quatro programas. Como o SAP-1 so soma e subtrai, fazemos "
    "multiplicacao com somas repetidas e divisao com subtracoes. Todos deram o resultado "
    "esperado, tanto na simulacao quanto na placa.",
    kicker="Validação", widths=[4.6,4.3,2.4])

content("Conclusão e trabalhos futuros",
    ["Implementamos um processador completo, do zero, em Verilog",
     "Entendemos a arquitetura (datapath, FSM, palavra de controle) e o funcionamento",
     "Verificado por simulação e validado na placa real",
     "Trabalhos futuros:",
     ("instrução de escrita (STA), desvios (JMP) para laços, rumo ao SAP-2",1)],
    "Pessoa 3",
    "Para concluir: construimos um processador inteiro e entendemos tanto a arquitetura "
    "quanto o funcionamento. Verificamos com testes e validamos na placa. Como evolucao, "
    "daria pra adicionar escrita na memoria e desvios, chegando ao SAP-2. Obrigado.",
    kicker="Encerramento",
    takeaway="Teoria + implementação + verificação + placa: o SAP-1 completo.")

closing()
finalize()

prs.save("SAP1_apresentacao_4partes.pptx")
print("PPTX gerado: SAP1_apresentacao_4partes.pptx (%d slides)"%len(prs.slides._sldIdLst))
